import os
from PyPDF2 import PdfReader
from pathlib import Path
import fitz
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import csv
from odf.opendocument import load
from odf.text import P, H
import zipfile
import xml.etree.ElementTree as ET
import re
from sentence_transformers import SentenceTransformer, util
from rapidfuzz import fuzz



def extraire_texte_odt(chemin_fichier):
    texte = ''
    try:
        with zipfile.ZipFile(chemin_fichier, 'r') as odt:
            xml_bytes = odt.read('content.xml')
            xml_str = xml_bytes.decode('utf-8')  # Convertit en texte brut
            # Enlève les balises XML pour récupérer tout le texte
            texte = re.sub(r'<[^>]+>', '', xml_str)
    except Exception as e:
        print(f"Erreur lors de l'extraction ODT de {chemin_fichier} : {e}")
    return texte.lower()

def extraire_texte_csv(chemin_fichier):
    texte = ''
    try:
        with open(chemin_fichier, newline='', encoding='utf-8') as fichier_csv:
            lecteur = csv.reader(fichier_csv)
            for ligne in lecteur:
                texte += ' '.join(ligne) + '\n'
    except Exception as e:
        print(f"Erreur lors de l'extraction CSV de {chemin_fichier} : {e}")
    return texte.lower()

def extraire_texte_pptx(chemin_fichier):
    prs = Presentation(chemin_fichier)
    texte = ''
    for diapositive in prs.slides:
        for forme in diapositive.shapes:
            if hasattr(forme, "text"):
                texte += forme.text + '\n'
    return texte.lower()

def extraire_texte_xlsx(chemin_fichier):
    wb = load_workbook(chemin_fichier, data_only=True)
    texte = ''
    for feuille in wb.worksheets:
        for ligne in feuille.iter_rows(values_only=True):
            texte += ' '.join([str(cellule) for cellule in ligne if cellule]) + '\n'
    return texte.lower()

def extraire_texte_docx(chemin_fichier):
    texte = ''
    try:
        doc = Document(chemin_fichier)
        for para in doc.paragraphs:
            if para.text.strip():  # on ignore les paragraphes vides
                texte += para.text + '\n'
    except Exception as e:
        print(f"Erreur lors de l'extraction DOCX de {chemin_fichier} : {e}")
    return texte.lower()

def extraire_texte_pdf(chemin_fichier):
    texte = ""
    try:
        doc = fitz.open(chemin_fichier)
        for page in doc:
            texte += page.get_text()
    except Exception as e:
        print(f"Erreur lors de l'extraction du texte de {chemin_fichier}: {e}")
    return texte.lower()


def extraire_texte(chemin_fichier):
    extension = Path(chemin_fichier).suffix.lower()
    try:
        if extension == ".pdf":
            return extraire_texte_pdf(chemin_fichier)
        elif extension == ".txt":
            with open(chemin_fichier, "r", encoding="utf-8") as f:
                return f.read().lower()
        elif extension == ".docx":
            return extraire_texte_docx(chemin_fichier)
        elif extension == ".xlsx":
            return extraire_texte_xlsx(chemin_fichier)
        elif extension == ".pptx":
            return extraire_texte_pptx(chemin_fichier)
        elif extension == ".csv":
            return extraire_texte_csv(chemin_fichier)
        elif extension == ".odt":
            return extraire_texte_odt(chemin_fichier)
    except Exception as e:
        print(f"⚠️ Erreur avec {chemin_fichier} : {e}")
    return ""

_modele_texte = None

def get_modele_texte():
    global _modele_texte
    if _modele_texte is None:
        _modele_texte = SentenceTransformer('all-MiniLM-L6-v2')
    return _modele_texte

def correspondance_semantique(texte, requete, seuil=0.5):
    try:
        modele = get_modele_texte()
        vecteur_texte = modele.encode(texte, convert_to_tensor=True)
        vecteur_requete = modele.encode(requete, convert_to_tensor=True)
        score = util.cos_sim(vecteur_texte, vecteur_requete)
        return score.item() > seuil
    except Exception as e:
        print(f"Erreur IA : {e}")
        return False

def correspondance_semantique_par_blocs(texte, requete, seuil=0.5, taille_bloc=500):
    try:
        modele = get_modele_texte()
        vecteur_requete = modele.encode(requete, convert_to_tensor=True)

        for i in range(0, len(texte), taille_bloc):
            bloc = texte[i:i+taille_bloc]
            vecteur_bloc = modele.encode(bloc, convert_to_tensor=True)
            score = util.cos_sim(vecteur_bloc, vecteur_requete)
            if score.item() > seuil:
                return True
        return False
    except Exception as e:
        print(f"Erreur IA par blocs : {e}")
        return False

    
def correspondance_fuzzy(texte, requete, seuil=80):
    lignes = texte.split('\n')
    for ligne in lignes:
        if fuzz.partial_ratio(ligne, requete) > seuil:
            return True
    return False


def recherche_mot_cle(dossier, mot_cle):
    mot_cle = mot_cle.lower()
    resultats = []

    for racine, _, fichiers in os.walk(dossier):
        for fichier in fichiers:
            chemin_fichier = os.path.join(racine, fichier)
            fichier_basename = fichier.lower()
            dossier_chemin = racine.lower()

            if mot_cle in fichier_basename or mot_cle in dossier_chemin:
                resultats.append(chemin_fichier)
                continue

            extension = os.path.splitext(fichier)[1].lower()
            texte = ''

            try:
                if extension == '.txt':
                    with open(chemin_fichier, 'r', encoding='utf-8', errors='ignore') as f:
                        texte = f.read().lower()

                elif extension == '.pdf':
                    texte = extraire_texte_pdf(chemin_fichier)

                elif extension == '.docx':
                    texte = extraire_texte_docx(chemin_fichier)

                elif extension == '.xlsx':
                    texte = extraire_texte_xlsx(chemin_fichier)

                elif extension == '.pptx':
                    texte = extraire_texte_pptx(chemin_fichier)

                elif extension == '.csv':
                    with open(chemin_fichier, 'r', encoding='utf-8', errors='ignore') as f:
                        texte = f.read().lower()

                elif extension == '.odt':
                    texte = extraire_texte_odt(chemin_fichier)

                if (
                    mot_cle in texte
                    or correspondance_fuzzy(texte, mot_cle)
                    or correspondance_semantique_par_blocs(texte, mot_cle)

                ):
                    resultats.append(chemin_fichier)

            except Exception as e:
                pass  

    return resultats


def rechercher_fichiers(dossier, extensions):
    fichiers_trouves = []
    for racine, _, fichiers in os.walk(dossier):
        for fichier in fichiers:
            if fichier.lower().endswith(extensions):
                fichiers_trouves.append(os.path.join(racine, fichier))
    return fichiers_trouves